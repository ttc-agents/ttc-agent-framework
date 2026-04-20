#!/usr/bin/env python3
"""
Knowledge Base Converter
Converts Word (.docx), PowerPoint (.pptx), and PDF (.pdf) files to plain text.
Uses a JSON index for delta processing - only converts new or modified files.

Usage:
    python3 convert_to_knowledge_base.py                                # default: all customers (legacy mode)
    python3 convert_to_knowledge_base.py --force                        # reconvert all files
    python3 convert_to_knowledge_base.py --source /path/to/folder       # specific source folder
    python3 convert_to_knowledge_base.py --source /path --dest KB/Finance  # custom KB destination
    python3 convert_to_knowledge_base.py --customer "ENOC"              # customer-scoped refresh
    python3 convert_to_knowledge_base.py --customer "ENOC" --force      # reconvert one customer
"""

import os
import sys
import json
import hashlib
import argparse
from pathlib import Path
from datetime import datetime

# ── Configuration ─────────────────────────────────────────────────────────────

SOURCE_ROOT = Path("/Users/joergpietzsch/Library/CloudStorage/OneDrive-TTCGlobal/Sales/Customer")
KB_ROOT     = Path("/Users/joergpietzsch/AI-Vault/Claude Folder/Knowledge Base")
INDEX_FILE  = KB_ROOT / "_index.json"
REGISTRY_FILE = KB_ROOT / "_customer_registry.json"
VENV_PYTHON = Path("/Users/joergpietzsch/AI-Vault/.venv/bin/python3")

SUPPORTED_EXTENSIONS = {".docx", ".pptx", ".pdf"}

# Directory names to skip anywhere in the tree
EXCLUDE_DIR_NAMES = {"AI-INFO - DO NOT DELETE", ".venv", "_vectordb", "__pycache__", ".git"}

# ── Helpers ───────────────────────────────────────────────────────────────────

def load_index(index_path: Path = INDEX_FILE) -> dict:
    if index_path.exists():
        with open(index_path, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}


def save_index(index: dict, index_path: Path = INDEX_FILE):
    index_path.parent.mkdir(parents=True, exist_ok=True)
    with open(index_path, "w", encoding="utf-8") as f:
        json.dump(index, f, indent=2, ensure_ascii=False)


def load_registry() -> dict:
    if REGISTRY_FILE.exists():
        with open(REGISTRY_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    return {"customers": {}}


def resolve_customer(name_or_slug: str) -> dict:
    """Return a customer registry entry by display name or slug. Raises KeyError if not found."""
    reg = load_registry()
    customers = reg.get("customers", {})
    # direct slug match
    if name_or_slug in customers:
        return customers[name_or_slug]
    # case-insensitive display name match
    for slug, entry in customers.items():
        if entry.get("display_name", "").lower() == name_or_slug.lower():
            return entry
    # slugified match
    import re
    slug = re.sub(r"[^a-z0-9]+", "-", name_or_slug.lower()).strip("-")
    if slug in customers:
        return customers[slug]
    raise KeyError(f"Customer not found in registry: {name_or_slug}")


def discover_files(source: Path) -> list:
    """Walk `source`, returning supported files while skipping EXCLUDE_DIR_NAMES."""
    found = []
    for dirpath, dirnames, filenames in os.walk(source):
        # Prune excluded directories in-place so os.walk skips them
        dirnames[:] = [d for d in dirnames if d not in EXCLUDE_DIR_NAMES]
        for fn in filenames:
            p = Path(dirpath) / fn
            if p.suffix.lower() in SUPPORTED_EXTENSIONS:
                found.append(p)
    return found


def file_checksum(path: Path) -> str:
    """MD5 checksum of a file for change detection."""
    h = hashlib.md5()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def relative_txt_path(source_path: Path, source_root: Path, kb_dest: Path) -> Path:
    """Mirror the source folder structure under kb_dest, with .txt extension."""
    rel = source_path.relative_to(source_root)
    return kb_dest / rel.with_suffix(".txt")


# ── Converters ────────────────────────────────────────────────────────────────

def convert_docx(src: Path) -> str:
    from docx import Document
    doc = Document(src)
    lines = []
    for para in doc.paragraphs:
        if para.text.strip():
            lines.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                lines.append(row_text)
    return "\n".join(lines)


def convert_pptx(src: Path) -> str:
    from pptx import Presentation
    prs = Presentation(src)
    lines = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        lines.append(f"--- Slide {slide_num} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


def convert_pdf(src: Path) -> str:
    try:
        import pdfminer.high_level as pdfminer
        return pdfminer.extract_text(str(src))
    except ImportError:
        pass
    # Fallback: use macOS built-in pdftotext if available
    import subprocess
    result = subprocess.run(["pdftotext", str(src), "-"], capture_output=True, text=True)
    if result.returncode == 0:
        return result.stdout
    raise RuntimeError(f"No PDF reader available. Install pdfminer.six: pip install pdfminer.six")


CONVERTERS = {
    ".docx": convert_docx,
    ".pptx": convert_pptx,
    ".pdf":  convert_pdf,
}


# ── Main logic ────────────────────────────────────────────────────────────────

def process_file(src: Path, dst: Path, index: dict, force: bool) -> str:
    """
    Returns: 'converted', 'skipped', or 'error:<message>'
    """
    src_key = str(src)
    checksum = file_checksum(src)
    mtime = os.path.getmtime(src)

    # Delta check
    if not force and src_key in index:
        if index[src_key].get("checksum") == checksum:
            return "skipped"

    try:
        converter = CONVERTERS[src.suffix.lower()]
        text = converter(src)
    except Exception as e:
        return f"error:{e}"

    # Write output mirroring the folder structure
    dst.parent.mkdir(parents=True, exist_ok=True)
    with open(dst, "w", encoding="utf-8") as f:
        f.write(text)

    # Update index entry
    index[src_key] = {
        "source":       str(src),
        "destination":  str(dst),
        "checksum":     checksum,
        "mtime":        mtime,
        "converted_at": datetime.now().isoformat(),
        "extension":    src.suffix.lower(),
    }
    return "converted"


def main():
    parser = argparse.ArgumentParser(description="Convert documents to Knowledge Base text files.")
    parser.add_argument("--force", action="store_true", help="Reconvert all files, ignoring the index.")
    parser.add_argument("--source", type=Path, default=None, help="Override source folder.")
    parser.add_argument("--dest",   type=Path, default=None, help="Override KB destination folder.")
    parser.add_argument("--customer", type=str, default=None,
                        help="Customer name or slug — looks up primary/AI-INFO folder from registry and routes there.")
    args = parser.parse_args()

    # Customer-scoped mode
    if args.customer:
        try:
            entry = resolve_customer(args.customer)
        except KeyError as e:
            print(f"Error: {e}", file=sys.stderr)
            print("Hint: run kb_bootstrap_customer.sh to register the customer first.", file=sys.stderr)
            sys.exit(2)
        source = Path(entry["primary_folder"])
        ai_info = Path(entry["ai_info_folder"])
        kb_dest = ai_info / "converted"
        index_path = ai_info / "_index.json"
        print(f"Customer      : {entry['display_name']}  (slug: {args.customer})")
        print(f"Primary folder: {source}")
    else:
        source = args.source if args.source else SOURCE_ROOT
        kb_dest = args.dest if args.dest else KB_ROOT
        index_path = INDEX_FILE

    index = load_index(index_path)

    # Discover all supported files (excluding AI-INFO folders etc.)
    all_files = discover_files(source)

    print(f"Found {len(all_files)} supported files in {source}")
    print(f"Destination KB folder : {kb_dest}")
    print(f"Index file            : {index_path}")
    print(f"Index contains {len(index)} previously converted entries\n")

    converted = 0
    skipped   = 0
    errors    = []

    for i, src in enumerate(all_files, start=1):
        dst = relative_txt_path(src, source, kb_dest)
        result = process_file(src, dst, index, force=args.force)

        if result == "converted":
            converted += 1
            print(f"[{i}/{len(all_files)}] CONVERTED  {src.relative_to(source)}")
        elif result == "skipped":
            skipped += 1
        else:
            error_msg = result.replace("error:", "")
            errors.append((src, error_msg))
            print(f"[{i}/{len(all_files)}] ERROR      {src.relative_to(source)} — {error_msg}")

        # Save index incrementally every 50 files
        if i % 50 == 0:
            save_index(index, index_path)

    # Final index save
    save_index(index, index_path)

    print(f"\n── Summary ───────────────────────────────")
    print(f"  Converted : {converted}")
    print(f"  Skipped   : {skipped} (unchanged)")
    print(f"  Errors    : {len(errors)}")
    print(f"  Index     : {index_path}")
    if errors:
        print(f"\n── Errors ────────────────────────────────")
        for src, msg in errors:
            print(f"  {src.name}: {msg}")


if __name__ == "__main__":
    main()
